from pptx import Presentation
import json

ppt = Presentation('d:/cricket/BPL.pptx')
print('slides', len(ppt.slides))

player_texts = []
for i, slide in enumerate(ppt.slides, start=1):
    texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]
    if texts:
        player_texts.append((i, texts))

print('slides_with_text', len(player_texts))
print('first 40 slide names:')
for idx, texts in player_texts[:40]:
    print(idx, texts)

all_names = [texts[0] for _, texts in player_texts if texts]
print('total_names', len(all_names))
print('unique_names', len(set(all_names)))

players = json.load(open('d:/cricket/players.json', 'r', encoding='utf-8'))
print('players_json', len(players))
print('first 40 players:')
for idx, p in enumerate(players[:40], start=1):
    print(idx, p.get('name'), p.get('registrationId'))
