"""
Import players from BPL.pptx: add missing entries, assign photos from slides
to players who lack a local /player_photos/ image.
"""
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parent
PPT = BASE_DIR / "BPL.pptx"
PLAYERS_JSON = BASE_DIR / "players.json"
OUT_DIR = BASE_DIR / "public" / "player_photos"
OUT_DIR.mkdir(parents=True, exist_ok=True)

NAME_ALIASES = {}


def norm_name(name):
    n = re.sub(r"^t+", "", (name or "").strip(), flags=re.I)
    n = re.sub(r"[^a-z0-9\s]", " ", n.lower())
    n = re.sub(r"\s+", " ", n).strip()
    return NAME_ALIASES.get(n, n)


def norm_roll(reg_id):
    return re.sub(r"[^a-zA-Z0-9]", "", (reg_id or "")).upper()


def parse_category_hand(lines):
    category = "Batsman"
    hand = "Right"
    for line in lines[1:]:
        low = line.lower().strip()
        if low in ("left", "right"):
            hand = line.strip().capitalize()
        elif "bowl" in low:
            category = "Bowler"
        elif "all" in low:
            category = "All-rounder"
        elif "bat" in low:
            category = "Batsman"
    return hand, category


def has_local_photo(player):
    url = player.get("photoUrl") or ""
    return url.startswith("/player_photos/")


def names_match(ppt_name, player_name):
    a = norm_name(ppt_name)
    b = norm_name(player_name)
    if a == b:
        return True
    at = a.split()
    bt = b.split()
    if len(at) == 1 or len(bt) == 1:
        return False
    return set(at) == set(bt)


def find_player(players, ppt_name, used_indices):
    for i, p in enumerate(players):
        if i in used_indices:
            continue
        if names_match(ppt_name, p["name"]):
            return i, p
    return None, None


def extract_slide_data(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())

    if not texts:
        return None

    name_line = texts[0]
    hand, category = parse_category_hand(texts)
    role = f"{hand} Hand {category}"

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    image_shape = None
    if pictures:
        # Prefer a portrait-shaped, reasonably sized picture (likely the player's headshot).
        imgs = []
        for s in pictures:
            w = getattr(s, "width", 0)
            h = getattr(s, "height", 0)
            area = int(w) * int(h)
            aspect = (h / w) if w else 0
            imgs.append((s, area, aspect))

        # Portrait candidates (height >= width * ~0.9)
        portraits = [t for t in imgs if t[2] >= 0.9]
        if portraits:
            max_area = max(a for (_, a, _) in imgs)
            # filter out extremely tiny icons
            portraits_filtered = [t for t in portraits if t[1] > max_area * 0.001]
            chosen = (min(portraits_filtered, key=lambda t: t[1])
                      if portraits_filtered else min(portraits, key=lambda t: t[1]))
            image_shape = chosen[0]
        else:
            # fallback: pick a reasonably small picture (not the huge template)
            max_area = max(a for (_, a, _) in imgs)
            candidates = [t for t in imgs if t[1] > max_area * 0.001 and t[1] < max_area * 0.9]
            if candidates:
                image_shape = min(candidates, key=lambda t: t[1])[0]
            else:
                # last resort, pick the smallest image
                image_shape = min(imgs, key=lambda t: t[1])[0]

    return {
        "name": re.sub(r"^t+", "", name_line, flags=re.I).strip(),
        "role": role,
        "hand": hand,
        "category": category,
        "image_shape": image_shape,
        "picture_shapes": pictures,
    }


def save_image(shape, dest: Path):
    try:
        img = shape.image
        blob = img.blob
        ext = img.ext or "png"
    except Exception:
        return None
    if not ext.startswith("."):
        ext = "." + ext
    dest = dest.with_suffix(ext)
    dest.write_bytes(blob)
    return dest


def main():
    prs = Presentation(str(PPT))
    players = json.load(PLAYERS_JSON.open(encoding="utf-8"))
    seen_rolls = {norm_roll(p.get("registrationId")) for p in players if norm_roll(p.get("registrationId"))}
    seen_phones = set()

    for p in players:
        phone = re.sub(r"\D", "", p.get("phone") or "")
        if phone:
            seen_phones.add(phone)

    ppt_players = []
    for slide_idx in range(1, len(prs.slides)):
        slide = prs.slides[slide_idx]
        data = extract_slide_data(slide)
        if data and data["name"]:
            ppt_players.append(data)

    added = []
    photos_updated = []

    for ppt in ppt_players:
        slug = re.sub(r"[^a-z0-9]", "", norm_name(ppt["name"]))[:20] or "player"
        reg = f"BPL_{slug.upper()}"
        n = 1
        while reg in seen_rolls:
            reg = f"BPL_{slug.upper()}_{n}"
            n += 1
        seen_rolls.add(reg)

        player = {
            "name": ppt["name"].title() if ppt["name"].islower() else ppt["name"],
            "role": ppt["role"],
            "class": "BPL",
            "registrationId": reg,
            "phone": "",
            "photoUrl": "",
        }
        players.append(player)
        idx = len(players) - 1
        added.append(player["name"])

        rid = norm_roll(player.get("registrationId")) or f"BPL_{idx}"
        dest_base = OUT_DIR / rid

        if ppt["image_shape"] is not None:
            saved = save_image(ppt["image_shape"], dest_base)
            # if extraction failed, try other picture shapes on the same slide
            if not saved:
                for other in ppt.get("picture_shapes", []):
                    if other is ppt["image_shape"]:
                        continue
                    saved = save_image(other, dest_base)
                    if saved:
                        break
            if saved:
                player["photoUrl"] = f"/player_photos/{saved.name}"
                photos_updated.append(f"{player['name']} ({player.get('registrationId')})")

    # Second pass: assign PPT photos to any still without local photo if name matches unused ppt entry
    for ppt in ppt_players:
        if ppt["image_shape"] is None:
            continue
        for i, player in enumerate(players):
            if has_local_photo(player):
                continue
            if names_match(ppt["name"], player["name"]):
                rid = norm_roll(player.get("registrationId")) or f"BPL_{i}"
                saved = save_image(ppt["image_shape"], OUT_DIR / rid)
                if not saved:
                    for other in ppt.get("picture_shapes", []):
                        if other is ppt["image_shape"]:
                            continue
                        saved = save_image(other, OUT_DIR / rid)
                        if saved:
                            break
                if saved:
                    player["photoUrl"] = f"/player_photos/{saved.name}"
                    label = f"{player['name']} ({player.get('registrationId')})"
                    if label not in photos_updated:
                        photos_updated.append(label)
                break

    json.dump(players, PLAYERS_JSON.open("w", encoding="utf-8"), indent=2, ensure_ascii=False)

    print(f"PPT players parsed: {len(ppt_players)}")
    print(f"Total in database: {len(players)}")
    if added:
        print(f"\nAdded {len(added)} new players:")
        for n in added:
            print(f"  + {n}")
    if photos_updated:
        print(f"\nPhotos added from PPT ({len(photos_updated)}):")
        for n in photos_updated:
            print(f"  * {n}")

    still_no_photo = [p["name"] for p in players if not has_local_photo(p)]
    if still_no_photo:
        print(f"\nStill no local photo ({len(still_no_photo)}):")
        for n in still_no_photo:
            print(f"  - {n}")


if __name__ == "__main__":
    main()
