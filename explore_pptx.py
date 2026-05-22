from pptx import Presentation
from pathlib import Path

PPT = Path(
    r"c:\Users\HP\AppData\Local\Packages\5319275A.WhatsAppDesktop_cv1g1gvanyjgm\LocalState\sessions\83FD23E5B3732A6F59229BBAF620A70C8CFA94CB\transfers\2026-21\BPL.pptx"
)

prs = Presentation(str(PPT))
print("slides", len(prs.slides))
for i, slide in enumerate(prs.slides):
    texts = []
    imgs = 0
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip().replace("\n", " | "))
        if shape.shape_type == 13:  # PICTURE
            imgs += 1
    print(f"\n--- Slide {i+1} (images={imgs}) ---")
    for t in texts:
        print(t[:200])
